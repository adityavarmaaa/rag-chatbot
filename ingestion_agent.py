# agents/ingestion_agent.py

import fitz  # PyMuPDF for PDFs
import docx
import pptx
import pandas as pd
import os

def parse_text(file_path):
    ext = os.path.splitext(file_path)[1].lower()

    try:
        if ext == ".pdf":
            doc = fitz.open(file_path)
            return "\n".join(page.get_text() for page in doc)

        elif ext == ".docx":
            doc = docx.Document(file_path)
            return "\n".join(p.text for p in doc.paragraphs)

        elif ext == ".pptx":
            pres = pptx.Presentation(file_path)
            return "\n".join(
                shape.text for slide in pres.slides for shape in slide.shapes if hasattr(shape, "text")
            )

        elif ext == ".csv":
            df = pd.read_csv(file_path)
            return df.to_string(index=False)

        elif ext in [".txt", ".md"]:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()

        else:
            return ""

    except Exception as e:
        print(f"[Error] Failed to parse {file_path}: {e}")
        return ""
